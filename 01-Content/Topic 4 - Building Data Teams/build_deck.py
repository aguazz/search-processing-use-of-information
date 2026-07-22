"""
Builds "Topic 4 - Building Data Teams (CUNEF style).pptx" from the approved
outline at `04-Course Redesign 2026-27/Slide Conversion Outlines/Building Data
Teams - Outline.md`, using the cunef-pptx skill's builder library.

Source content: the outline above, itself mapped from the original
`Topic 4 - Building Data Teams.pptx` (source slide numbers cited per section
in the outline; renamed from "Topic 2" to "Topic 4" to match the corrected
2026-27 numbering in Course_Timetable_2026-27.xlsx). No speaker notes, per
this conversion project's standing decision (Slide_Conversion_Plan.md) —
these are lectures already delivered.
"""

import sys
import os
from pptx.util import Inches

_THIS_DIR = os.path.dirname(os.path.abspath(__file__))
sys.path.insert(0, os.path.join(_THIS_DIR, "..", "..", ".claude", "skills", "cunef-pptx"))
import cunef_pptx_builder as cb

prs = cb.new_presentation()

# ---------------------------------------------------------------------------
# Icon placement helper — icons are used ONLY on two-column comparison slides
# (the "page 7" type: one icon centered under each column), per the reviewed
# decision that icons read well there but clutter the 3-card / single-concept
# / citation slides. The 3-card icon helper was removed for that reason.
# Geometry mirrors add_two_column_slide's own box math (see that function);
# the icon sits in the empty space below each column's (short) bullet text.
# ---------------------------------------------------------------------------

_COL_LEFTS = [0.6, 6.8665]
_COL_WIDTH = 5.8665
_COL_ICON_TOP = Inches(5.2)
_COL_ICON_SIZE = Inches(0.55)


def add_column_icons(slide, icon_names):
    """icon_names: 2 icon names (or None to skip that column), left to right."""
    for i, name in enumerate(icon_names):
        if name is None:
            continue
        left = Inches(_COL_LEFTS[i] + (_COL_WIDTH - 0.55) / 2)
        cb.add_icon(slide, name, left, _COL_ICON_TOP, _COL_ICON_SIZE)

# ===========================================================================
# Title + Roadmap
# ===========================================================================

s = cb.add_title_slide(
    prs, "Building Data Teams",
    "Teamwork, Kanban & Your Group Project",
)

s = cb.add_roadmap_slide(prs, "Roadmap", [
    ("1", "Teamwork"),
    ("2", "Kanban"),
    ("3", "Project"),
    ("4", "Tools"),
])

# ===========================================================================
# Section: Teamwork & Project Management
# ===========================================================================

s = cb.add_section_divider(prs, 1, "Teamwork")

s = cb.add_blank_slide(prs, "Teamwork can be challenging", "1 | Teamwork")
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.3), prs.slide_width - Inches(1.2), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Most people think working in groups is easy. It isn't - not even "
          "for us, coordinating this course! Let's do an exercise that "
          "proves the point.")
cb.format_body_paragraph(p, 20)

s = cb.add_three_card_slide(
    prs, "The Marshmallow Challenge", "1 | Teamwork",
    [
        ("Objective", "Build the tallest freestanding structure, measured from the table to the top of the marshmallow."),
        ("Materials & rules", [
            "Use only what you're given: spaghetti, string, tape",
            "The marshmallow must sit on top",
            "Destroying or eating it disqualifies your team",
        ]),
        ("Time & prize", "15 minutes, timed against a countdown video. Winning team gets a prize!"),
    ],
)

s = cb.add_bullet_list_slide(
    prs, "Some questions for you", "1 | Teamwork",
    [
        "Was there a leader on your team?",
        "Who decided who the leader would be?",
        "If you had no leader, would designating one have helped?",
        "Did you have a plan?",
    ],
)

s = cb.add_two_column_slide(
    prs, "Some lessons: who performs how", "1 | Teamwork",
    "Perform poorly",
    [
        "Business-school teams often do worst",
        "They spend most of their time agreeing on ONE plan",
        "They run out of time rushing that plan - and it collapses",
    ],
    "Perform well",
    [
        "Kindergarteners are surprisingly good at this",
        "They skip planning and just start building",
        "Many attempts collapse, but each one teaches them something",
    ],
)
add_column_icons(s, ["triangle-alert", "circle-check"])

s = cb.add_three_card_slide(
    prs, "More lessons", "1 | Teamwork",
    [
        ("Prototype early", "The exercise rewards iterative design over one big upfront plan."),
        ("It's heavier than it looks", "The marshmallow is heavier than most people assume - budget for it structurally, not just visually."),
        ("Beware hidden assumptions", "Teams that add the marshmallow last, assuming it's light, often watch their structure collapse. Every project has its own hidden 'marshmallow.'"),
    ],
)

s = cb.add_statement_slide(
    prs,
    "Everything exciting at a company comes from projects - and every "
    "project succeeds or fails as a group, not as individuals. You will "
    "never work alone.",
)

# ===========================================================================
# Section: Kanban
# ===========================================================================

s = cb.add_section_divider(prs, 2, "Kanban")

s = cb.add_blank_slide(prs, "What is Kanban?", "2 | Kanban")
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.3), prs.slide_width - Inches(1.2), Inches(2.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Kanban is a visual way to manage work as it flows through a "
          "process: every task is visible on a board, you cap how much is "
          "in progress at once, and you keep improving the flow. Originally "
          "a factory-floor method, today it runs everything from software "
          "teams to your group project.")
cb.format_body_paragraph(p, 20)
cb.add_help_link(prs, s, "Kanban, explained (Atlassian)",
                 "https://www.atlassian.com/agile/kanban",
                 lead="Want the full guide? ")

s = cb.add_process_chain_slide(
    prs, "Where Kanban came from", "2 | Kanban",
    [
        ("1950s", "Taiichi Ohno conceives 'produce only what's needed, when needed'"),
        ("Piggly Wiggly", "Inspired by a US supermarket restocking just enough of each product"),
        ("Kanban cards", "Paper cards signal and track demand on the factory floor"),
        ("JIT", "Becomes the basis of Toyota's Just-In-Time methodology"),
    ],
)

# Merged: the old "Core values" (3 pillars) + "Kanban values in action"
# (7-item list) collapsed into one slide — the 3 core values, each with the
# concrete practices that express it drawn from the old 7-item list. Generic
# umbrella items ("Enable Lean/Agile culture") dropped rather than force-fit.
s = cb.add_three_card_slide(
    prs, "Core values of Kanban", "2 | Kanban",
    [
        ("Flow", [
            "Continuously and predictably deliver value",
            "Deliver often, in small increments",
        ]),
        ("Visibility", [
            "Make work and its workflow visible, so issues surface early",
            "Prioritize openly, so the team sees what matters most",
        ]),
        ("Limit WIP", [
            "Cap how much work is in progress at once",
            "Less multitasking - shorter cycle times, higher quality",
        ]),
    ],
)

s = cb.add_bullet_list_slide(
    prs, "Limit Work in Progress", "2 | Kanban",
    [
        "Reduces cycle time per task",
        "Ensures WIP is the highest-priority task",
        "Reduces or eliminates queues",
        "Reduces multi-tasking by team members",
        "Rule of thumb: keep WIP at 5 or fewer items per person",
    ],
)

s = cb.add_process_chain_slide(
    prs, "Kanban boards", "2 | Kanban",
    [
        ("To do", "Work waiting to start"),
        ("Doing", "Work in progress"),
        ("Done", "Work completed"),
    ],
)

s = cb.add_three_card_slide(
    prs, "It's all about the cards", "2 | Kanban",
    [
        ("One card, one task", "Every work item is a separate card, tracking progress visually as it moves through the workflow."),
        ("What a card shows", [
            "Who's responsible",
            "A short description of the task",
            "How long it's estimated to take",
        ]),
        ("The goal: reach Done", "Each card moves from To Do through every column until it lands in Done - once every card is there, the project is finished."),
    ],
)

s = cb.add_statement_slide(
    prs,
    "You won't nail your board's setup on the first try - you'll add or "
    "remove columns, adjust cards, and rearrange priorities as you go. "
    "(SpaceX didn't land its first rocket booster either.)",
)

# ===========================================================================
# Section: Group Project
# ===========================================================================

s = cb.add_section_divider(prs, 3, "Group Project")

s = cb.add_bullet_list_slide(
    prs, "Your deliverables", "3 | Group Project",
    [
        "Written Report (citations tracked via Zotero)",
        "Video Presentation, submitted in Canvas",
        "Kanban Board in Trello",
        "The Zotero library you used for the project",
        "Full detail: Canvas -> Group Project module",
    ],
)

s = cb.add_three_card_slide(
    prs, "The Written Report", "3 | Group Project",
    [
        ("Methodology", "Explain neural networks to a client who knows nothing about them - no complicated math."),
        ("Application", "Real-world examples of neural networks applied to your assigned industry."),
        ("References", "Every source you used, tracked in Zotero."),
    ],
)

s = cb.add_blank_slide(prs, "The Video Presentation", "3 | Group Project")
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.3), prs.slide_width - Inches(1.2), Inches(2.5))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("Submit a video presentation of your report - maximum 10 "
          "minutes. Use animation tools, record yourselves, whatever "
          "works. Be creative - try to impress the client.")
cb.format_body_paragraph(p, 20)

s = cb.add_three_card_slide(
    prs, "Building your Kanban board", "3 | Group Project",
    [
        ("Define your flow", "At least 3 columns - To Do, Doing, Done. Add more if your work needs them (e.g. an Editing column)."),
        ("Set a WIP limit", "Rule of thumb: 1.5x the number of people who can work a step at once. Doesn't apply to To Do or Done."),
        ("Split big tasks", "Anything over a couple of days, or needing multiple people, becomes several smaller cards you can track independently."),
    ],
)

s = cb.add_bullet_list_slide(
    prs, "Grading", "3 | Group Project",
    [
        "Description of Neural Networks methodology - 2 pts",
        "Application to the assigned industry - 2 pts",
        "Use of the Kanban Board - 2 pts",
        "Presentation (instructor-assessed) - 3 pts",
        "Correct use of references (Zotero) - 1 pt",
        "Total - 10 pts",
    ],
)

s = cb.add_two_column_slide(
    prs, "Due dates & group meetings", "3 | Group Project",
    "Due dates",
    [
        "See Canvas for the exact submission deadline",
    ],
    "Group meetings",
    [
        "Ask any time - don't leave it to the last minute",
        "I'll check in with a different group each week",
    ],
)
add_column_icons(s, ["calendar", "message-circle"])

# ===========================================================================
# Section: Trello & Zotero
# ===========================================================================

s = cb.add_section_divider(prs, 4, "Trello & Zotero")

s = cb.add_blank_slide(prs, "Trello setup", "4 | Trello & Zotero")
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.3), prs.slide_width - Inches(1.2), Inches(2.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("You'll share a Trello board with me as your virtual Kanban "
          "board - your project's To Do / Doing / Done, online. A short "
          "setup screencast is on the way; until then, the official "
          "Trello guide below walks you through creating your first board.")
cb.format_body_paragraph(p, 20)
cb.add_help_link(prs, s, "Trello Guide: set up your first board (official)",
                 "https://trello.com/guide",
                 lead="New to Trello? ")

s = cb.add_blank_slide(prs, "Zotero setup", "4 | Trello & Zotero")
tb = s.shapes.add_textbox(Inches(0.6), Inches(2.3), prs.slide_width - Inches(1.2), Inches(2.8))
tf = tb.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.text = ("You'll track your Group Project's citations and references "
          "here, synced with MS Word so your bibliography builds itself. "
          "A short setup screencast is on the way; until then, the "
          "official Zotero guide below covers the essentials.")
cb.format_body_paragraph(p, 20)
cb.add_help_link(prs, s, "Zotero Quick Start Guide (official)",
                 "https://www.zotero.org/support/quick_start_guide",
                 lead="New to Zotero? ")

# ===========================================================================
# Closing
# ===========================================================================

s = cb.add_statement_slide(
    prs,
    "This week: how to run a Kanban board, manage it with Trello, and "
    "track your references with Zotero - the exact toolkit you'll use "
    "for your Group Project all semester.",
)

s = cb.add_closing_slide(prs)

# ===========================================================================

out_pptx = os.path.join(_THIS_DIR, "Topic 4 - Building Data Teams (CUNEF style).pptx")
out_pdf = os.path.join(_THIS_DIR, "Topic 4 - Building Data Teams (CUNEF style).pdf")
cb.save(prs, out_pptx)
print("Saved:", out_pptx, "-", len(prs.slides), "slides")
cb.export_to_pdf(out_pptx, out_pdf)
print("Exported:", out_pdf)
